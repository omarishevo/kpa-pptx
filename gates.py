import streamlit as st
from pptx import Presentation
from PIL import Image
from io import BytesIO
import os

# Title
st.title("🚦 Vehicle Traffic and Congestion at KPA Gates – Slide Viewer")

# Load PowerPoint file
pptx_path = "OMARPPTX.pptx"  # Make sure this file is in the same directory or use file uploader
prs = Presentation(pptx_path)

# Create a list of slides
slides = list(prs.slides)
total_slides = len(slides)

# Select slide
slide_number = st.slider("Select Slide Number", 1, total_slides, 1)
slide = slides[slide_number - 1]

st.header(f"📊 Slide {slide_number} of {total_slides}")

# Extract text from the selected slide
slide_text = ""
for shape in slide.shapes:
    if hasattr(shape, "text"):
        slide_text += shape.text + "\n"

st.markdown("### 📝 Slide Content")
st.text(slide_text.strip())

# Extract and display images (if any)
st.markdown("### 🖼️ Visuals/Charts")
for shape in slide.shapes:
    if shape.shape_type == 13:  # Picture type
        image = shape.image
        image_bytes = image.blob
        img = Image.open(BytesIO(image_bytes))
        st.image(img, use_column_width=True)

# Optional: Upload your own .pptx
st.sidebar.markdown("## 📤 Upload a Different Presentation")
uploaded_file = st.sidebar.file_uploader("Upload .pptx File", type=["pptx"])
if uploaded_file:
    with open("uploaded_pptx.pptx", "wb") as f:
        f.write(uploaded_file.getbuffer())
    st.sidebar.success("File uploaded successfully. Reload the app with the new file.")
